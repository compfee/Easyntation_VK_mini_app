# example
from pptx.util import Inches

from pptx_generator import PPTGenerator
from summarization import summ_text
from text_parsing import parse_docx
import pickle
import os

if __name__ == '__main__':

  # with open('summ_text.txt') as f:
  #   text = f.read()
  doc_loc = os.path.dirname(os.path.dirname(__file__)) +"/doc2.docx"
  parse_docx(doc_loc)

  with open(os.path.dirname(os.path.dirname(__file__)) + '/data/result.pickle', "rb") as f:
    article_text_full = pickle.load(f)

  print(article_text_full)
  ppt = PPTGenerator('Example PPT')
  ppt.add_slide(ppt.get_slide_layout(0))
  ppt.add_title(title_text=article_text_full[0]['Heading'], index_slide=0)

  ppt.add_slide(ppt.get_slide_layout(1))
  ppt.add_text(article_text_full[1]['Heading'], summ_text(article_text_full[1]), index_slide=1)
  ppt.add_image('Image', '1.jpg', index_slide=1)

  # print(article_text_full[8]['Enum'][0])
  # print(article_text_full[8]['Enum'][3])
  list_text = [article_text_full[8]['Enum'][0], article_text_full[8]['Enum'][2],
               article_text_full[8]['Enum'][3], article_text_full[8]['Enum'][4],
               article_text_full[8]['Enum'][5]]
  ppt.add_slide(ppt.get_slide_layout(1))
  print(list_text)
  ppt.add_text(article_text_full[8]['Subtitle'], list_text, index_slide=2, location='center')


  ppt.add_slide(ppt.get_slide_layout(1))
  ppt.add_text(article_text_full[16]['Heading'], summ_text(article_text_full[16]), index_slide=3)
  ppt.add_image('Image', '2.jpg', index_slide=3)

  # ppt.add_slide(ppt.get_slide_layout(1))
  # ppt.add_text(article_text_full[22]['Heading'], summ_text(article_text_full[22]), index_slide=4)

  # ppt.add_slide(ppt.get_slide_layout(1))
  # ppt.add_text(article_text_full[16]['Heading'], summ_text(article_text_full[16]), index_slide=3)
  #
  # ppt.add_slide(ppt.get_slide_layout(1))
  # ppt.add_text(article_text_full[16]['Heading'], summ_text(article_text_full[16]), index_slide=3)

  ppt.save_ppt()